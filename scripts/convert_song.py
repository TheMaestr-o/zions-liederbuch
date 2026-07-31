#!/usr/bin/env python3
"""
End-to-end pptx -> Keynote converter for one song.

Usage:
    python3 convert_song.py 145
    python3 convert_song.py 145 --no-verify

Finds "{number} - *.pptx" in this folder, parses it (title, verses, refrain),
builds the new-design .key via make_song.py, then reads the result back
through Keynote and asserts it matches what was asked for. Prints one of:
    OK   <number> <title>  (<n> verses, refrain: yes/no)
    FAIL <reason>
so a whole batch can be driven without inspecting slide dumps by hand.
"""

import glob
import os
import re
import subprocess
import sys
import time

from pptx import Presentation

from make_song import build_song, osa, as_str, KEYNOTE_ID

SONGS_DIR = "/Users/ohnedan/Developer/zions-liederbuch/songs"
PPTX_SEARCH_DIRS = [
    "/Users/ohnedan/Downloads/Zions Liederbuch/0-100",
    "/Users/ohnedan/Downloads/Zions Liederbuch/100-200",
    "/Users/ohnedan/Downloads/Zions Liederbuch/200-300",
    "/Users/ohnedan/Downloads/Zions Liederbuch/300-400",
    "/Users/ohnedan/Downloads/Zions Liederbuch/400-500",
    "/Users/ohnedan/Downloads/Zions Liederbuch/500-600",
    "/Users/ohnedan/Downloads/Zions Liederbuch/600-705",
]


def find_pptx(number: str) -> str:
    matches = []
    for d in PPTX_SEARCH_DIRS:
        matches += glob.glob(f"{d}/{number} - *.pptx") or glob.glob(f"{d}/{number}.pptx")
    if not matches:
        raise FileNotFoundError(f"no pptx found for song {number}")
    if len(matches) > 1:
        raise RuntimeError(f"ambiguous pptx matches for {number}: {matches}")
    return matches[0]


def clean_block(text: str) -> str:
    text = text.replace("\x0b", "\n")
    lines = [ln.strip() for ln in text.split("\n")]
    lines = [ln for ln in lines if ln != ""]
    return "\n".join(lines)


def strip_leading_label(text: str, patterns) -> str:
    lines = text.split("\n")
    if lines:
        lines[0] = re.sub(patterns, "", lines[0]).strip()
    return "\n".join(lines)


def parse_pptx(path: str):
    pres = Presentation(path)
    slides = []
    for slide in pres.slides:
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slides.append(texts)

    # drop leading/trailing blank slides
    while slides and not slides[0]:
        slides.pop(0)
    while slides and not slides[-1]:
        slides.pop()

    if not slides:
        raise ValueError(f"can't find cover slide (number+title) in {path}")

    cover = slides.pop(0)
    # Some covers (e.g. song 36) put "Lied Nr. N" and the title as two LINES of a
    # SINGLE text box instead of two separate shapes. Split it so the number/title
    # detection below sees them as two blocks, exactly as in the usual two-shape case.
    if len(cover) == 1:
        parts = [p for p in clean_block(cover[0]).split("\n") if p]
        if len(parts) >= 2:
            cover = [parts[0], "\n".join(parts[1:])]
    if len(cover) < 2:
        raise ValueError(f"can't find cover slide (number+title) in {path}")

    number = None
    number_idx = None
    for idx, block in enumerate(cover):
        joined = " ".join(block.replace("\x0b", " ").split())
        # "Lie+d" tolerates source typos like "Lieed Nr.637" (song 637)
        m = re.search(r"Lie+d\s*Nr\.?\s*(\d+)", joined, re.IGNORECASE) or re.search(r"^(\d+)$", joined)
        if m:
            number = m.group(1)
            number_idx = idx
            break
    if number is None:
        # Last resort: the filename is authoritative anyway (convert() overrides with
        # it), so only the title block still has to be identified. Any block carrying
        # some "Nr. <digits>" label is the number block, the other one is the title.
        m = re.search(r"(\d+)", os.path.basename(path))
        if m:
            number = m.group(1)
            for idx, block in enumerate(cover):
                joined = " ".join(block.replace("\x0b", " ").split())
                if re.search(r"Nr\.?\s*\d+", joined, re.IGNORECASE):
                    number_idx = idx
                    break
    if number is None:
        raise ValueError(f"no song number found on cover slide: {cover}")
    title_block = next((b for i, b in enumerate(cover) if i != number_idx), cover[0])
    title = " ".join(clean_block(title_block).split("\n"))

    content_blocks = []
    for texts in slides:
        if not texts:
            continue
        content_blocks.append(clean_block(texts[0]))

    is_refrain = [bool(re.match(r"^(Refr?\.?:?|Chor:?)", b, re.IGNORECASE)) for b in content_blocks]

    verses = []
    refrain = None
    for block, refr in zip(content_blocks, is_refrain):
        if refr:
            block = strip_leading_label(block, r"^(Refr?\.?:?|Chor:?)\s*")
            if refrain is None:
                refrain = block
        else:
            block = strip_leading_label(block, r"^\d+\.\s*")
            verses.append(block)

    return dict(number=number, title=title, verses=verses, refrain=refrain)


def restart_keynote():
    try:
        osa(f'tell application id "{KEYNOTE_ID}" to quit saving no')
    except Exception:
        pass
    subprocess.run(["pkill", "-9", "-f", "Keynote Creator Studio"], check=False)
    time.sleep(3)


def verify(output_path: str, number: str, title: str, verses: list, refrain: str) -> list:
    """Read the finished deck back and compare it with what was asked for.

    This Keynote build sometimes stops answering AppleEvents right after a save -
    even `get name of documents` then returns -1712, while the app sits idle at 3%
    CPU (seen on song 226). The saved file itself is fine, so a wedged app must not
    condemn the song: force-restart Keynote and read it back once more. Only the
    failure path is affected, so songs that verify first time behave exactly as before.
    """
    try:
        return _verify_once(output_path, number, title, verses, refrain)
    except RuntimeError as e:
        detail = str(e).strip().splitlines()[-1][:110]
        print(f"   WARN verify failed ({detail}); restarting Keynote and retrying once")
        restart_keynote()
        return _verify_once(output_path, number, title, verses, refrain)


def _verify_once(output_path: str, number: str, title: str, verses: list, refrain: str) -> list:
    problems = []
    osa(f'''
    with timeout of 180 seconds
        tell application id "{KEYNOTE_ID}"
            open POSIX file {as_str(output_path)}
            delay 2
        end tell
    end timeout
    tell application "System Events" to set visible of application process "Keynote" to false
    ''')
    # \x01 stands in for embedded newlines in multi-line verse/refrain text so the
    # dump stays exactly one real line per slide, and \x02 separates the fields.
    # The field separator must NOT be a character that can occur in hymn text: "|"
    # used to be used here, but German repeat marks "|: ... :|" appear in many songs
    # (137, 163, 269, 359, 385, 509, 512, 636, 697) and split their verses into
    # bogus extra fields, which made verification fail on perfectly good decks.
    dump = osa(f'''
    on flatten(t)
        set AppleScript's text item delimiters to linefeed
        set theItems to text items of t
        set AppleScript's text item delimiters to (ASCII character 1)
        set t2 to theItems as string
        set AppleScript's text item delimiters to ""
        return t2
    end flatten

    with timeout of 180 seconds
        tell application id "{KEYNOTE_ID}"
            set sep to (ASCII character 2)
            set theDoc to front document
            set n to count of slides of theDoc
            set out to (n as string) & linefeed
            repeat with i from 1 to n
                set t to text items of (slide i of theDoc)
                set out to out & i & sep & (count of t)
                repeat with it_ in t
                    set out to out & sep & my flatten(object text of it_)
                end repeat
                set out to out & linefeed
            end repeat
            return out
        end tell
    end timeout
    ''')
    osa(f'''
    with timeout of 180 seconds
        tell application id "{KEYNOTE_ID}" to close front document saving no
    end timeout
    ''')

    lines = dump.split("\n")
    n = int(lines[0])
    expected = 3 + len(verses) * (2 if refrain else 1)
    if n != expected:
        problems.append(f"slide count {n}, expected {expected}")

    slide_lines = lines[1:1 + n]
    slides = []
    for line in slide_lines:
        parts = line.split("\x02")
        slides.append([p.replace("\x01", "\n") for p in parts[2:]])

    cover = slides[1] if len(slides) > 1 else []
    if len(cover) < 3 or cover[0] != number or cover[1] != title or cover[2] != title:
        problems.append(f"cover mismatch: {cover}")

    pos = 1
    for i, verse_text in enumerate(verses, start=1):
        pos += 1
        if pos >= len(slides):
            problems.append(f"missing verse slide {i}")
            break
        s = slides[pos]
        if len(s) < 5 or s[0] != str(i) or s[1] != verse_text or s[2] != number or s[4] != verse_text:
            problems.append(f"verse {i} mismatch at slide {pos + 1}: {s}")
        if refrain:
            pos += 1
            if pos >= len(slides):
                problems.append(f"missing refrain slide after verse {i}")
                break
            r = slides[pos]
            if len(r) < 5 or r[1] != refrain or r[2] != number or r[4] != refrain:
                problems.append(f"refrain mismatch after verse {i} at slide {pos + 1}: {r}")

    return problems


def convert(number: str, do_verify: bool = True):
    pptx_path = find_pptx(number)
    song = parse_pptx(pptx_path)
    if song["number"] != number:
        print(f"WARN filename number {number} != cover slide number {song['number']}, using {number}")
        song["number"] = number

    safe_title = re.sub(r'[/:]', '-', song["title"])
    output_path = f'{SONGS_DIR}/{number} - {safe_title}.key'
    build_song(
        number=song["number"],
        title=song["title"],
        verses=song["verses"],
        refrain=song["refrain"],
        output_path=output_path,
    )

    status = "OK"
    problems = []
    if do_verify:
        problems = verify(output_path, song["number"], song["title"], song["verses"], song["refrain"])
        status = "OK" if not problems else "FAIL"

    print(f'{status} {number} "{song["title"]}" ({len(song["verses"])} verses, refrain: {"yes" if song["refrain"] else "no"})')
    for p in problems:
        print("   -", p)
    return status == "OK"


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: convert_song.py <number> [--no-verify]")
        sys.exit(1)
    num = sys.argv[1]
    do_verify = "--no-verify" not in sys.argv
    ok = convert(num, do_verify)
    sys.exit(0 if ok else 1)
